# -*- coding: utf-8 -*-
"""Generate the Regnify investor pitch deck template, matching the brand template.
v2: slide titles moved into the header band (level with the wordmark, above the
faint divider) to reclaim vertical space; adds AI-engine + scalability slides;
updated team slide (Andrew, Willie, Luan)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
CHARCOAL   = RGBColor(0x1A,0x19,0x17)
WARM_DARK  = RGBColor(0x2E,0x2B,0x25)
GOLD       = RGBColor(0xC8,0xA8,0x4B)
LIGHT_GOLD = RGBColor(0xE8,0xD9,0x9A)
IVORY      = RGBColor(0xF0,0xE8,0xD8)
CREAM      = RGBColor(0xF5,0xF1,0xEA)
MID_WARM   = RGBColor(0x9A,0x8E,0x7E)
DEEP_WARM  = RGBColor(0x4A,0x47,0x40)
CARD_BORDER= RGBColor(0xDE,0xDA,0xD2)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
FOOT_LIGHT = RGBColor(0xB0,0xA8,0x98)
FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def set_bg(slide, color):
    cSld = slide._element.find(qn('p:cSld'))
    for b in cSld.findall(qn('p:bg')):
        cSld.remove(b)
    bg = cSld.makeelement(qn('p:bg'), {})
    bgPr = bg.makeelement(qn('p:bgPr'), {})
    fill = bg.makeelement(qn('a:solidFill'), {})
    clr = bg.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0],color[1],color[2])})
    fill.append(clr); bgPr.append(fill)
    bgPr.append(bg.makeelement(qn('a:effectLst'), {}))
    bg.append(bgPr)
    cSld.insert(0, bg)

def rect(slide, l, t, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w if line_w else 0.75)
    sp.shadow.inherit = False
    return sp

def txt(slide, l, t, w, h, runs, size=12, color=DEEP_WARM, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(para, str):
            segs = [(para, {})]
        elif isinstance(para, tuple):
            segs = [para]
        else:
            segs = [(s, {}) if isinstance(s, str) else s for s in para]
        for seg_text, ov in segs:
            r = p.add_run(); r.text = seg_text
            f = r.font
            f.name = ov.get('font', font)
            f.size = Pt(ov.get('size', size))
            f.bold = ov.get('bold', bold)
            f.italic = ov.get('italic', italic)
            f.color.rgb = ov.get('color', color)
    return tb

def bullets(slide, l, t, w, h, items, size=12, color=DEEP_WARM, gap=4, spacing=1.0,
            marker="•  ", marker_color=GOLD, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing; p.space_after = Pt(gap)
        ital = False; col = color
        if isinstance(it, tuple):
            it, opt = it
            ital = opt.get('italic', False); col = opt.get('color', color)
        if marker:
            rm = p.add_run(); rm.text = marker
            rm.font.name = FONT; rm.font.size = Pt(size); rm.font.color.rgb = marker_color; rm.font.bold = True
        r = p.add_run(); r.text = it
        r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = col; r.font.italic = ital
    return tb

PAGE = [0]
def chrome(slide, dark=False, footer=True, page=True):
    PAGE[0] += 1
    wm_color   = IVORY if dark else CHARCOAL
    div_color  = RGBColor(0x3E,0x3B,0x35) if dark else CARD_BORDER
    foot_color = DEEP_WARM if dark else FOOT_LIGHT
    rect(slide, 0, 0, 10, 0.12, fill=GOLD)
    rect(slide, 0, 0.82, 10, 0.0, line=div_color, line_w=0.75)
    txt(slide, 6.92, 0.18, 2.80, 0.42, "Regnify", size=22, color=wm_color, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, 7.08, 0.62, 2.64, 0.04, fill=GOLD)
    rect(slide, 9.76, 0.58, 0.09, 0.09, fill=GOLD, shape=MSO_SHAPE.OVAL)
    rect(slide, 9.88, 0.59, 0.07, 0.07, fill=LIGHT_GOLD, shape=MSO_SHAPE.OVAL)
    if footer:
        txt(slide, 0.42, 5.18, 4.0, 0.30, "Focus on work that matters.", size=9, color=foot_color)
    if page:
        txt(slide, 8.80, 5.18, 0.80, 0.30, "%02d" % PAGE[0], size=9, color=foot_color, align=PP_ALIGN.RIGHT)

def new(dark=False, bg=WHITE, footer=True, page=True):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, bg)
    chrome(s, dark=dark, footer=footer, page=page)
    return s

def head_title(slide, title, sub=None):
    """Slide title in the header band (left, level with the wordmark, above the line)."""
    txt(slide, 0.42, 0.17, 6.30, 0.48, title, size=22, color=CHARCOAL, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        txt(slide, 0.42, 0.94, 9.16, 0.40, sub, size=12.5, color=MID_WARM, italic=True)

def card(slide, l, t, w, h, header=None, body=None, header_color=CHARCOAL,
         body_items=None, body_size=11):
    rect(slide, l, t, w, h, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    rect(slide, l, t, w, 0.07, fill=GOLD)
    if header:
        txt(slide, l+0.18, t+0.20, w-0.36, 0.38, header, size=14, color=header_color, bold=True)
    if body:
        txt(slide, l+0.18, t+0.64, w-0.36, h-0.82, body, size=body_size, color=DEEP_WARM, spacing=1.05)
    if body_items:
        bullets(slide, l+0.18, t+0.62, w-0.36, h-0.80, body_items, size=body_size, gap=3, spacing=1.0)

def stat(slide, l, t, w, h, big, label, big_size=26):
    rect(slide, l, t, w, h, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    rect(slide, l, t, w, 0.07, fill=GOLD)
    txt(slide, l+0.18, t+0.20, w-0.36, 0.5, big, size=big_size, color=CHARCOAL, bold=True)
    txt(slide, l+0.18, t+0.74, w-0.36, h-0.88, label, size=10, color=MID_WARM, spacing=1.0)

def avatar(slide, path, l, t, d, imgw, imgh, ring=GOLD):
    """Circular headshot: square center-crop (biased up for portraits) + ellipse mask."""
    pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(d), Inches(d))
    if imgw > imgh:
        extra = (imgw - imgh) / imgw
        pic.crop_left = extra/2; pic.crop_right = extra/2
    elif imgh > imgw:
        extra = (imgh - imgw) / imgh
        pic.crop_top = extra*0.30; pic.crop_bottom = extra*0.70
    spPr = pic._element.spPr
    for tag in ('a:prstGeom','a:custGeom'):
        e = spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    geom = spPr.makeelement(qn('a:prstGeom'), {'prst':'ellipse'})
    geom.append(spPr.makeelement(qn('a:avLst'), {}))
    spPr.append(geom)
    if ring:
        rect(slide, l, t, d, d, fill=None, line=ring, line_w=1.5, shape=MSO_SHAPE.OVAL)
    return pic

GUIDE = {'italic': True, 'color': GOLD}

# ===================================================================
# 1 — COVER (dark)
# ===================================================================
s = new(dark=True, bg=CHARCOAL, footer=False, page=False)
txt(s, 0.42, 1.75, 9.16, 1.0, "Regnify", size=58, color=IVORY, bold=True, align=PP_ALIGN.CENTER)
rect(s, 2.6, 2.92, 4.8, 0.04, fill=GOLD)
txt(s, 1.0, 3.12, 8.0, 0.5,
    "AI-powered Licensed Representative lifecycle management for MAS-regulated firms",
    size=15, color=MID_WARM, align=PP_ALIGN.CENTER, italic=True)
txt(s, 1.0, 4.55, 8.0, 0.35, "Investor Pitch Deck", size=13, color=IVORY, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.0, 4.92, 8.0, 0.3,
    [[("Seed Round  ·  ", {'color':MID_WARM}), ("Confidential", {'color':GOLD}),
      ("  ·  ", {'color':MID_WARM}), ("[Month] 2026", {'color':MID_WARM, 'italic':True})]],
    size=11, align=PP_ALIGN.CENTER)

# ===================================================================
# 2 — PROBLEM
# ===================================================================
s = new()
head_title(s, "The problem", "Managing Licensed Representatives under MAS is manual, fragmented and high-risk.")
txt(s, 0.42, 1.50, 4.55, 2.7,
    [[("Today, FA firms and fund managers run the representative lifecycle across spreadsheets, "
       "email and disconnected vendors.", {})],
     [(" ", {'size':6})],
     [("Form 3A/B/C/D filings, Fit & Proper checks, CMFAS/CPD tracking and RNF submissions are "
       "stitched together by hand — with no single source of truth and no audit trail when MAS asks.", {})]],
    size=13, color=DEEP_WARM, spacing=1.1)
stat(s, 5.28, 1.45, 2.10, 1.40, "[xx] hrs", "manual effort per representative, per year")
stat(s, 7.48, 1.45, 2.10, 1.40, "[xx]+", "disconnected systems & vendors in a typical firm")
stat(s, 5.28, 2.98, 2.10, 1.40, "$[xx]k", "cost of a single Fit & Proper / RNF lapse")
stat(s, 7.48, 2.98, 2.10, 1.40, "Days", "to assemble evidence for a MAS inspection")
txt(s, 5.28, 4.50, 4.3, 0.3, "Fill: replace bracketed figures with your validated data points.",
    size=9, color=GOLD, italic=True)

# ===================================================================
# 3 — SOLUTION
# ===================================================================
s = new()
head_title(s, "The solution", "One platform for the entire Licensed Representative lifecycle — onboarding to exit.")
y=1.45; h=2.95
card(s, 0.42, y, 2.98, h, header="End-to-end lifecycle",
     body="Pre-onboarding, post-onboarding and cessation in one workflow — Form 3A/B/C/D, "
          "Fit & Proper, CMFAS mapping, CPD and RNF-ready data.")
card(s, 3.51, y, 2.98, h, header="Rep Passport",
     body="A portable representative record — employment history, CMFAS modules, CPD and a "
          "downloadable profile that moves with the representative between firms.")
card(s, 6.60, y, 2.98, h, header="AI-native compliance",
     body="AI agents advise, check workflows, draft documentation and ingest reports — "
          "turning regulatory know-how into an always-on layer across the platform.")

# ===================================================================
# 4 — PRODUCT SCOPE (ecosystem diagram, full bleed)
# ===================================================================
s = new(footer=False, page=True)
head_title(s, "Product scope")
txt(s, 0.42, 0.98, 5.93, 0.24, "REPRESENTATIVE LIFECYCLE MANAGEMENT  ·  Financial Institution",
    size=9.5, color=MID_WARM, bold=True)
txt(s, 6.55, 0.98, 3.03, 0.24, "REPRESENTATIVE HISTORY  ·  Representative",
    size=9.5, color=MID_WARM, bold=True)
stages = [("Pre-Onboarding",0.42),("Post-Onboarding",2.39),("Cessation / Exit",4.36)]
for name,lx in stages:
    rect(s, lx, 1.24, 2.05, 0.40, fill=GOLD, shape=MSO_SHAPE.PENTAGON)
    txt(s, lx, 1.24, 1.86, 0.40, name, size=11, color=CHARCOAL, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 6.55, 1.24, 3.03, 0.40, fill=CHARCOAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 6.55, 1.24, 3.03, 0.40, "Rep Passport", size=11, color=IVORY, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
def feat(l, w, items):
    rect(s, l, 1.74, w, 1.66, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    bullets(s, l+0.12, 1.86, w-0.24, 1.48, items, size=8.5, gap=2.5, spacing=1.0, marker="– ", marker_color=GOLD)
feat(0.42, 1.91, ["Form 3A/B/C/D workflow","Fit & Proper checks","CMFAS mapping by role",
                  "Background & reference checks","RNF submission-ready data","Audit trail"])
feat(2.39, 1.91, ["Regulated-activity tracking","Update of particulars","CPD compliance tracking",
                  "Activity attestation","Breach / incident logging","Recertification"])
feat(4.36, 1.91, ["Change / cessation tracking","Clean exit documentation","Hand-off evidence pack"])
feat(6.55, 3.03, ["Employment history","CMFAS modules & CPD records","Downloadable profile",
                  "Export / import across employers"])
rect(s, 0.42, 3.58, 9.16, 0.66, fill=WARM_DARK)
txt(s, 0.60, 3.66, 2.1, 0.5, [["REGNIFY"],[("AI-powered RegTech", {'size':8,'color':MID_WARM})]],
    size=11, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
agents = ["AI Regulatory Advisor","AI Workflow Checker","AI Documentation Agent",
          "AI Exam Cert Analyser","AI Report Ingester"]
ax = 2.75; aw = 1.32
for a in agents:
    rect(s, ax, 3.74, aw, 0.34, fill=CHARCOAL, line=GOLD, line_w=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, ax, 3.74, aw, 0.34, a, size=7.5, color=IVORY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ax += aw + 0.05
partners = ["Background Screening","L&D Partners","CMFAS Partners","AML/CFT/PF Partners","MAS","IBF"]
px = 0.42; pw = 1.49
for p in partners:
    rect(s, px, 4.42, pw, 0.50, fill=CREAM, line=CARD_BORDER, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, px, 4.42, pw, 0.50, p, size=8, color=DEEP_WARM, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    px += pw + 0.05
txt(s, 0.42, 4.98, 9.16, 0.25, "Partner & regulator ecosystem", size=8.5, color=MID_WARM, italic=True, align=PP_ALIGN.CENTER)

# ===================================================================
# 5 — THE AI ENGINE
# ===================================================================
s = new()
head_title(s, "The AI engine", "AI-enhanced workflows and smart documents — with governance built for regulators.")
# left: capabilities
txt(s, 0.42, 1.42, 5.0, 0.3, "AI capabilities", size=14, color=CHARCOAL, bold=True)
caps = ["Document data extraction","Certification gap analysis","Regulatory advisory chatbot",
        "Pre-approval check pipeline","Form pre-fill from documents","Enhanced agentic chatbot",
        "3rd-party report comparison","Alert & notification system","CPD tracking & recommendations",
        "Regulatory change detection"]
cols_x = [0.42, 3.02]; cw = 2.50
for i, c in enumerate(caps):
    col = i % 2; rowi = i // 2
    cx = cols_x[col]; cy = 1.82 + rowi * 0.535
    rect(s, cx, cy, cw, 0.46, fill=CREAM, line=CARD_BORDER, line_w=0.6, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx+0.14, cy, 0.28, 0.46, "%d" % (i+1), size=10, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx+0.44, cy, cw-0.56, 0.46, c, size=9, color=CHARCOAL, anchor=MSO_ANCHOR.MIDDLE, spacing=0.95)
# right: governance panel
rect(s, 5.74, 1.42, 3.84, 3.30, fill=WARM_DARK)
txt(s, 5.96, 1.58, 3.4, 0.3, "Governance principles", size=14, color=GOLD, bold=True)
gov = [("Human-in-the-loop","Humans make the final call on every decision."),
       ("Citation-grounded","Every answer cites a specific MAS source."),
       ("Full audit trail","Every action logged for regulatory inspection."),
       ("Privacy by design","PII redacted before any external LLM call."),
       ("Cost-optimised routing","The right model tier for each task.")]
gy = 2.04
for hd, bd in gov:
    txt(s, 5.96, gy, 3.46, 0.26, hd, size=11, color=IVORY, bold=True)
    txt(s, 5.96, gy+0.26, 3.46, 0.26, bd, size=9, color=MID_WARM, spacing=0.95)
    gy += 0.55
txt(s, 0.42, 4.80, 5.0, 0.3, "Granular capabilities — keep for NDA'd investors; the public deck stays high-level.",
    size=8.5, color=GOLD, italic=True)

# ===================================================================
# 6 — WHY NOW
# ===================================================================
s = new()
head_title(s, "Why now", "Three forces make this the moment for AI-native RegTech in Singapore.")
y=1.45; h=2.95
card(s, 0.42, y, 2.98, h, header="Regulation tightening",
     body="MAS continues to raise the bar on representative conduct, Fit & Proper and "
          "accountability — increasing the compliance load on every licensed firm.")
card(s, 3.51, y, 2.98, h, header="Rising cost of compliance",
     body="Headcount and manual processes don't scale. Firms need to do more oversight with "
          "the same team, and absorb fewer errors.")
card(s, 6.60, y, 2.98, h, header="AI is finally ready",
     body="Modern AI can read regulation, map obligations and draft documentation reliably — "
          "making an AI-native compliance layer viable for the first time.")
txt(s, 0.42, 4.55, 9.16, 0.3, "Fill: add a specific MAS regulatory milestone or date that sharpens the timing.",
    size=9, color=GOLD, italic=True)

# ===================================================================
# 7 — MARKET OPPORTUNITY
# ===================================================================
s = new()
head_title(s, "Market opportunity", "Built bottom-up from Singapore, expanding across ASEAN.")
stat(s, 0.42, 1.50, 2.85, 2.05, "$[TAM]", "TAM — all MAS-regulated firms managing representatives, plus ASEAN expansion")
stat(s, 3.51, 1.50, 2.85, 2.05, "$[SAM]", "SAM — Singapore FA firms & fund managers in your initial segment")
stat(s, 6.60, 1.50, 2.85, 2.05, "$[SOM]", "SOM — realistic 3-year capture at target win rate")
txt(s, 0.42, 3.80, 9.16, 0.6,
    [[("Bottom-up logic:  ", {'bold':True, 'color':CHARCOAL}),
      ("# of Licensed Representatives / firms in segment  ×  annual contract value  ×  attach rate. "
       "Show the calculation — investors trust bottom-up over top-down.", {'color':GOLD, 'italic':True})]],
    size=11, spacing=1.1)

# ===================================================================
# 8 — SCALABILITY / EXPANSION (NEW)
# ===================================================================
s = new()
head_title(s, "Scalability", "The same engine extends to any regulated profession, in any regulated market.")
# 2x2 growth matrix
bl = (2.50, 3.02); br = (6.00, 3.02); tl = (2.50, 1.52); tr = (6.00, 1.52)
cw2, ch2 = 3.35, 1.40
def qcell(pos, t1, t2, gold=False, emphasis=False):
    l,t = pos
    rect(s, l, t, cw2, ch2, fill=(GOLD if gold else CREAM),
         line=(GOLD if emphasis else CARD_BORDER), line_w=(1.5 if emphasis else 0.75))
    txt(s, l+0.18, t+0.22, cw2-0.36, 0.5, t1, size=12.5,
        color=(CHARCOAL if gold else CHARCOAL), bold=True, spacing=0.95)
    txt(s, l+0.18, t+0.72, cw2-0.36, 0.55, t2, size=9.5,
        color=(DEEP_WARM if gold else DEEP_WARM), spacing=0.98)
qcell(tl, "Financial services × ASEAN", "Same product, new markets: Malaysia, Indonesia, wider ASEAN.")
qcell(tr, "New verticals × ASEAN", "The full regulated-professions opportunity, region-wide.", emphasis=True)
qcell(bl, "Financial services × Singapore", "Licensed Representatives — LIVE today (beachhead).", gold=True)
qcell(br, "New verticals × Singapore", "Real estate, healthcare, legal & other regulated professions.")
# axis cues
txt(s, 2.50, 1.20, 5.0, 0.26, "↑  Expand across ASEAN markets", size=10, color=GOLD, bold=True)
txt(s, 2.50, 4.50, 6.85, 0.28, "Add regulated verticals  →", size=10, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
txt(s, 0.42, 4.86, 9.16, 0.3,
    "One engine — lifecycle, credentials & AI compliance — powers every cell. Each expansion reuses the same core.",
    size=10, color=MID_WARM, italic=True, align=PP_ALIGN.CENTER)

# ===================================================================
# 9 — COMPETITION & MOAT
# ===================================================================
s = new()
head_title(s, "Competition & moat", "Why Regnify wins, and why it's hard to copy.")
rows = [
    ["", "Manual / in-house", "Generic GRC tools", "Regnify"],
    ["Built for MAS LR lifecycle", "No", "Partial", "Yes — purpose-built"],
    ["Rep Passport / portability", "No", "No", "Yes"],
    ["AI-native compliance layer", "No", "Bolt-on", "Core"],
    ["Partner & regulator ecosystem", "No", "Limited", "Embedded"],
]
tl_, tt_ = 0.42, 1.50
colw = [2.1, 1.3, 1.3, 1.3]; rh = 0.46
yy = tt_
for ri, row in enumerate(rows):
    xx = tl_
    for ci, cell in enumerate(row):
        is_head = ri == 0; is_reg = ci == 3
        fill = CHARCOAL if is_head else (CREAM if is_reg else WHITE)
        rect(s, xx, yy, colw[ci], rh, fill=fill, line=CARD_BORDER, line_w=0.5)
        if is_reg and not is_head:
            rect(s, xx, yy, colw[ci], 0.04, fill=GOLD)
        col = IVORY if is_head else (CHARCOAL if (is_reg or ci==0) else DEEP_WARM)
        bold = is_head or ci==0 or is_reg
        txt(s, xx+0.08, yy, colw[ci]-0.12, rh, cell, size=9.5, color=col, bold=bold,
            align=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER), anchor=MSO_ANCHOR.MIDDLE)
        xx += colw[ci]
    yy += rh
rect(s, 6.72, 1.50, 2.86, 2.30, fill=WARM_DARK)
txt(s, 6.92, 1.66, 2.5, 0.35, "Our moat", size=14, color=GOLD, bold=True)
bullets(s, 6.92, 2.12, 2.5, 1.6,
        ["Regulatory depth in MAS rules","Proprietary lifecycle + Rep Passport data","Partner & regulator integrations","Switching costs once embedded"],
        size=10, color=IVORY, gap=6, marker="•  ", marker_color=GOLD)

# ===================================================================
# 10 — BUSINESS MODEL
# ===================================================================
s = new()
head_title(s, "Business model", "Recurring SaaS, priced to the value of compliance assurance.")
y=1.45; h=2.95
card(s, 0.42, y, 2.98, h, header="Revenue model",
     body_items=["SaaS subscription (annual)","Per-firm platform fee","Per-representative seat tier","Partner / marketplace revenue (future)"], body_size=11)
card(s, 3.51, y, 2.98, h, header="Pricing",
     body_items=["Land with core lifecycle module","Expand with AI agents & Rep Passport","Tiered by firm size / rep count",("Fill: target ACV per firm", GUIDE)], body_size=11)
card(s, 6.60, y, 2.98, h, header="Unit economics",
     body_items=[("Fill: CAC", GUIDE),("Fill: LTV & LTV:CAC", GUIDE),("Fill: gross margin %", GUIDE),("Fill: payback period", GUIDE)], body_size=11)

# ===================================================================
# 11 — GO-TO-MARKET
# ===================================================================
s = new()
head_title(s, "Go-to-market", "A partner-led wedge into MAS-regulated firms, then land-and-expand.")
steps = [("01","Target","Start with [segment] — e.g. mid-size FA firms & fund managers with [n] representatives."),
         ("02","Channel","Direct founder-led sales + warm intros via compliance and industry networks."),
         ("03","Partners","Co-sell through L&D, CMFAS, screening and AML partners; credibility via IBF / MAS engagement."),
         ("04","Expand","Land with one module, expand to full lifecycle + AI agents across the firm.")]
x=0.42; w=2.24
for num,hh,b in steps:
    rect(s, x, 1.45, w, 2.9, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    rect(s, x, 1.45, w, 0.07, fill=GOLD)
    txt(s, x+0.16, 1.62, w-0.3, 0.6, num, size=24, color=GOLD, bold=True)
    txt(s, x+0.16, 2.24, w-0.3, 0.35, hh, size=13, color=CHARCOAL, bold=True)
    is_guide = "[" in b
    txt(s, x+0.16, 2.66, w-0.3, 1.5, b, size=10.5, color=(GOLD if is_guide else DEEP_WARM),
        italic=is_guide, spacing=1.05)
    x += w + 0.13
txt(s, 0.42, 4.52, 9.16, 0.3, "Fill: replace [segment] / [n] with your beachhead definition and target metrics.",
    size=9, color=GOLD, italic=True)

# ===================================================================
# 12 — TRACTION
# ===================================================================
s = new()
head_title(s, "Traction", "Early signal that the market wants this.")
stat(s, 0.42, 1.50, 2.10, 1.55, "[x]", "design partners / pilots")
stat(s, 2.62, 1.50, 2.10, 1.55, "[x]", "LOIs or signed firms")
stat(s, 4.82, 1.50, 2.10, 1.55, "[x]", "representatives in pipeline")
stat(s, 7.02, 1.50, 2.10, 1.55, "[x]", "partner agreements")
txt(s, 0.42, 3.30, 9.16, 0.35, "Proof points", size=14, color=CHARCOAL, bold=True)
bullets(s, 0.42, 3.72, 9.16, 1.1,
        [("Fill: name pilots / design partners and what they validated", GUIDE),
         ("Fill: quote a customer or regulatory-engagement milestone", GUIDE),
         ("Fill: waitlist, MoUs, or revenue to date", GUIDE)],
        size=11, gap=5)

# ===================================================================
# 13 — TEAM & FOUNDERS (updated)
# ===================================================================
s = new()
head_title(s, "Team & founders", "Two ex-MAS founders, backed by award-winning AI engineering.")
txt(s, 0.42, 1.32, 6.0, 0.24, "CO-FOUNDERS", size=10, color=GOLD, bold=True)
txt(s, 6.54, 1.32, 2.93, 0.24, "CORE TEAM", size=10, color=GOLD, bold=True)
team = [
 (0.42, "assets/andrew.jpg", (1537,1849), "Andrew Quake", "Co-Founder",
   ["~30 yrs scaling APAC businesses in global MNCs (ex-CEO / GM)",
    "Domain: payments, customer experience & contact centres",
    "Started career at MAS · MAS Scholar",
    "MBA, London Business School · Accountancy (Hons), NTU"]),
 (3.48, "assets/willie.png", (712,762), "Willie Tan", "Co-Founder",
   ["Seasoned compliance leader — advisory, surveillance & regulatory engagement",
    "Builds digital-first frameworks for cross-border buy-side FIs",
    "Started career at MAS (6 yrs)",
    "MSc Fin Eng (NUS) · MA Digital Mgmt · Accountancy (Hons), NTU"]),
 (6.54, "assets/luan.jpg", (730,874), "Minh Luan Pham", "Founding Head of Technology",
   ["12+ yrs enterprise software · 5 yrs production AI / ML",
    "Finalist, VnExpress AI Awards 2024 (anti-hallucination GenAI)",
    "Top 35 FINTECH 2025 — AI analytics live at KAFI Securities",
    "Certs: ML in Finance (NYU) · Stanford ML · TensorFlow · NLP"]),
]
cw3 = 2.93
for l, img, (iw,ih), name, role, bio in team:
    rect(s, l, 1.55, cw3, 3.40, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    rect(s, l, 1.55, cw3, 0.07, fill=GOLD)
    avatar(s, img, l+0.18, 1.74, 0.92, iw, ih)
    txt(s, l+1.24, 1.80, cw3-1.42, 0.34, name, size=12.5, color=CHARCOAL, bold=True, spacing=0.95)
    txt(s, l+1.24, 2.14, cw3-1.42, 0.5, role, size=9.5, color=GOLD, italic=True, spacing=0.92)
    bullets(s, l+0.18, 2.84, cw3-0.36, 2.05, bio, size=8.5, gap=3.5, spacing=0.98, marker="– ", marker_color=GOLD)

# ===================================================================
# 14 — FINANCIAL PROJECTIONS
# ===================================================================
s = new()
head_title(s, "Financial projections", "Three-year P&L — fill with your model's outputs.")
fin_rows = [
    ["S$'000", "Year 1", "Year 2", "Year 3"],
    ["Customers (firms)", "", "", ""],
    ["Revenue (ARR)", "", "", ""],
    ["Gross profit", "", "", ""],
    ["Operating expenses", "", "", ""],
    ["EBITDA", "", "", ""],
    ["Cash balance / runway", "", "", ""],
]
tl_, tt_ = 0.42, 1.45; colw=[3.3,1.95,1.95,1.95]; rh=0.45
yy=tt_
for ri,row in enumerate(fin_rows):
    xx=tl_
    for ci,cell in enumerate(row):
        head = ri==0
        fill = CHARCOAL if head else (CREAM if ri%2==0 else WHITE)
        rect(s, xx, yy, colw[ci], rh, fill=fill, line=CARD_BORDER, line_w=0.5)
        disp = cell
        if not head and ci>0 and cell=="":
            disp = "—"
        col = IVORY if head else (CHARCOAL if ci==0 else MID_WARM)
        txt(s, xx+0.1, yy, colw[ci]-0.16, rh, disp, size=10.5, color=col, bold=(head or ci==0),
            align=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.RIGHT), anchor=MSO_ANCHOR.MIDDLE)
        xx+=colw[ci]
    yy+=rh
txt(s, 0.42, yy+0.10, 9.16, 0.3,
    "Fill: drop in the figures from your financial model; keep key assumptions in an appendix slide.",
    size=9, color=GOLD, italic=True)

# ===================================================================
# 15 — RISKS & CHALLENGES
# ===================================================================
s = new()
head_title(s, "Risks & challenges", "What could go wrong — and how we de-risk it.")
risks = [
    ("Regulatory change", "MAS rules evolve and could shift requirements.", "AI monitors changes; modular rules engine adapts fast."),
    ("Long sales cycles", "Regulated firms buy slowly and cautiously.", "Partner-led intros + pilots shorten time-to-trust."),
    ("Data security & trust", "We handle sensitive representative data.", "Security-first architecture; align to MAS expectations."),
    ("Incumbent response", "GRC vendors could move down-market.", "Depth, Rep Passport data and ecosystem create lock-in."),
]
yy=1.50; rhh=0.66
rect(s, 0.42, yy, 2.3, rhh-0.24, fill=CHARCOAL); rect(s,2.72,yy,3.0,rhh-0.24,fill=CHARCOAL); rect(s,5.72,yy,3.86,rhh-0.24,fill=CHARCOAL)
txt(s,0.5,yy,2.2,rhh-0.24,"Risk",size=10,color=IVORY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,2.8,yy,2.9,rhh-0.24,"Why it matters",size=10,color=IVORY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,5.8,yy,3.7,rhh-0.24,"Mitigation",size=10,color=IVORY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
yy+=rhh-0.24
for i,(r,w,m) in enumerate(risks):
    fill = CREAM if i%2==0 else WHITE
    rect(s,0.42,yy,2.3,rhh,fill=fill,line=CARD_BORDER,line_w=0.5)
    rect(s,2.72,yy,3.0,rhh,fill=fill,line=CARD_BORDER,line_w=0.5)
    rect(s,5.72,yy,3.86,rhh,fill=fill,line=CARD_BORDER,line_w=0.5)
    txt(s,0.5,yy,2.18,rhh,r,size=10,color=CHARCOAL,bold=True,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0)
    txt(s,2.8,yy,2.86,rhh,w,size=9.5,color=DEEP_WARM,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0)
    txt(s,5.8,yy,3.7,rhh,m,size=9.5,color=DEEP_WARM,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0)
    yy+=rhh

# ===================================================================
# 16 — ROADMAP
# ===================================================================
s = new()
head_title(s, "Roadmap", "Milestones over the next 18-24 months.")
rect(s, 0.7, 3.05, 8.6, 0.03, fill=GOLD)
phases = [("Now","Pilot & design partners","Validate core lifecycle module"),
          ("+6 mo","First paying firms","Launch Rep Passport + AI agents"),
          ("+12 mo","Scale in Singapore","Land-and-expand; partner channel live"),
          ("+18 mo","ASEAN expansion","Second market entry")]
xs=[0.9, 3.0, 5.1, 7.2]
for i,(t,hh,b) in enumerate(phases):
    x=xs[i]
    rect(s, x-0.06, 2.95, 0.22, 0.22, fill=GOLD, shape=MSO_SHAPE.OVAL)
    above = i%2==0
    cy = 1.60 if above else 3.22
    rect(s, x+0.02, cy, 1.95, 1.18, fill=CREAM, line=CARD_BORDER, line_w=0.75)
    rect(s, x+0.02, cy, 1.95, 0.06, fill=GOLD)
    txt(s, x+0.14, cy+0.11, 1.7, 0.3, t, size=12, color=GOLD, bold=True)
    txt(s, x+0.14, cy+0.42, 1.72, 0.35, hh, size=11, color=CHARCOAL, bold=True, spacing=0.95)
    txt(s, x+0.14, cy+0.76, 1.72, 0.4, b, size=9, color=DEEP_WARM, spacing=0.95)
txt(s, 0.42, 4.78, 9.16, 0.3, "Fill: tune milestones and dates to your plan.", size=9, color=GOLD, italic=True)

# ===================================================================
# 17 — THE ASK
# ===================================================================
s = new()
head_title(s, "The ask", "What we're raising, and what it unlocks.")
rect(s, 0.42, 1.50, 4.0, 2.7, fill=WARM_DARK)
txt(s, 0.62, 1.70, 3.6, 0.3, "Raising", size=12, color=MID_WARM)
txt(s, 0.62, 2.00, 3.6, 0.7, "S$[X.X]M", size=40, color=GOLD, bold=True)
txt(s, 0.62, 2.85, 3.6, 0.3, "Seed round  ·  [equity / SAFE]", size=11, color=IVORY, italic=True)
txt(s, 0.62, 3.25, 3.6, 0.7, "Gives ~[xx] months runway to hit [key milestone].",
    size=11, color=MID_WARM, italic=True, spacing=1.05)
txt(s, 4.7, 1.50, 4.9, 0.35, "Use of funds", size=14, color=CHARCOAL, bold=True)
uses = [("Product & engineering","[xx]%"),("Compliance & regulatory","[xx]%"),
        ("Go-to-market & partnerships","[xx]%"),("Operations & runway","[xx]%")]
yy=2.05
for label,pct in uses:
    txt(s, 4.7, yy, 3.6, 0.3, label, size=11, color=DEEP_WARM)
    txt(s, 8.3, yy, 1.3, 0.3, pct, size=11, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
    rect(s, 4.7, yy+0.32, 4.9, 0.02, fill=CARD_BORDER)
    yy+=0.55

# ===================================================================
# 18 — CLOSING (dark)
# ===================================================================
s = new(dark=True, bg=CHARCOAL, footer=False, page=False)
txt(s, 0.5, 1.85, 9.0, 1.0, "Regnify", size=52, color=IVORY, bold=True, align=PP_ALIGN.CENTER)
rect(s, 2.6, 2.95, 4.8, 0.04, fill=GOLD)
txt(s, 1.0, 3.15, 8.0, 0.5, "Let's build the compliance backbone for MAS-regulated firms.",
    size=15, color=MID_WARM, align=PP_ALIGN.CENTER, italic=True)
txt(s, 1.0, 4.5, 8.0, 0.35,
    [[("ask@regnify.sg", {'color':IVORY,'bold':True}), ("    ·    ", {'color':MID_WARM}),
      ("regnify.sg", {'color':IVORY,'bold':True})]],
    size=13, align=PP_ALIGN.CENTER)

import os, sys
out = sys.argv[1] if len(sys.argv) > 1 else "Regnify Pitch Deck Template.pptx"
prs.save(out)
print("Saved", out)
